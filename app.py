from flask import Flask, request, jsonify
from langchain_community.document_loaders.pdf import PyPDFLoader
from langchain_community.vectorstores import FAISS
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.chains.combine_documents import create_stuff_documents_chain
from langchain_core.prompts import ChatPromptTemplate
from langchain.chains import create_retrieval_chain
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_core.runnables import RunnablePassthrough
from langchain_core.output_parsers import StrOutputParser
from langchain.prompts import ChatPromptTemplate, PromptTemplate
from langchain_community.document_loaders import PyPDFDirectoryLoader
from dotenv import load_dotenv
from flask_cors import CORS, cross_origin
from pptx import Presentation
from pptx.util import Inches
import google.generativeai as genai
import os

load_dotenv()

embeddings = HuggingFaceEmbeddings(model_name="mixedbread-ai/mxbai-embed-large-v1", encode_kwargs={'precision': 'binary'})

# Vector Store
faiss_path = os.getenv('FAISS_PATH', 'OSAnd')  # Use environment variable or default to OSAnd
try:
    one_bit_vectorstore = FAISS.load_local(faiss_path, embeddings, allow_dangerous_deserialization=True)
    retriever = one_bit_vectorstore.as_retriever(search_kwargs={"k": 10})
except Exception as e:
    print(f"Error loading FAISS index: {str(e)}")
    # Initialize empty FAISS index if loading fails
    one_bit_vectorstore = None
    retriever = None

app = Flask(__name__)
cors = CORS(app)
app.config['CORS_HEADERS'] = 'Content-Type'

# Configure the Gemini API
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
genai.configure(api_key=GOOGLE_API_KEY)

def create_slide(presentation, title, content_points):
    slide_layout = presentation.slide_layouts[1]  # Use the layout with title and content
    slide = presentation.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]
    
    title_placeholder.text = title
    
    for point in content_points:
        p = content_placeholder.text_frame.add_paragraph()
        p.text = point
        p.level = 0  # Bullet point level

def generate_presentation(slide_data, template_path, output_path):
    presentation = Presentation(template_path)
    
    for slide_info in slide_data:
        title = slide_info['title']
        content_points = slide_info['content']
        
        while content_points:
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            title_placeholder = slide.shapes.title
            content_placeholder = slide.placeholders[1]
            
            title_placeholder.text = title
            
            # Add points to the slide until it is full
            remaining_points = []
            for point in content_points:
                p = content_placeholder.text_frame.add_paragraph()
                p.text = point
                p.level = 0  # Bullet point level
                
                # Check if the content exceeds the slide's capacity
                if content_placeholder.text_frame.fit_text():
                    remaining_points.append(point)
                    content_placeholder.text_frame.text = content_placeholder.text_frame.text.rsplit('\n', 1)[0]
                    break
            
            content_points = remaining_points
    
    presentation.save(output_path)

def is_comparative_or_exploratory_question(question):
    """
    Detect if the question is comparative or exploratory
    """
    comparative_keywords = [
        'difference', 'compare', 'contrast', 'best', 'better', 'vs', 'versus', 
        'what is', 'explain', 'describe', 'how', 'why', 'characteristics', 
        'features', 'advantages', 'disadvantages'
    ]
    
    # Convert question to lowercase for case-insensitive matching
    lower_question = question.lower()
    
    # Check if any comparative keywords are in the question
    return any(keyword in lower_question for keyword in comparative_keywords)

def enhance_query(query):
    """Enhance the query to get more detailed responses"""
    if query.lower().startswith("what is "):
        enhanced_query = query.lower().replace("what is ", "Explain the concept of ", 1)
    elif query.lower().startswith("define "):
        enhanced_query = query.lower().replace("define ", "Provide a detailed definition of ", 1)
    elif query.lower().startswith("who is "):
        enhanced_query = query.lower().replace("who is ", "Explain the significance of ", 1)
    else:
        # Add a request for detail
        enhanced_query = f"Provide a comprehensive explanation of: {query}. Include examples and relevant context."

    return enhanced_query

@app.route("/", methods=['POST'])
def hello_world():
    content = request.json

    if retriever is None:
        return {"answer": "The vector store is not initialized. Please contact the administrator."}, 500

    # Gemini model configuration
    model = genai.GenerativeModel('gemini-1.5-pro')

    # Add basic validation for question content
    if 'question' not in content or not content['question'].strip():
        return {"answer": "Please provide a valid question."}, 400

    # Enhance the query
    enhanced_question = enhance_query(content['question'])

    # Improved template with flexible answering strategy
    template = """
    Context: {context}
    
    Question: {question}
    
Guidelines for Response:

1. Responses will be strictly based on the information available in the vector database. If the context provides direct information relevant to the question, the response will use that information verbatim, citing the exact document name and section as the source.

2. For comparative or exploratory questions, responses will be structured with a comprehensive explanation based only on relevant documents in the vector database. Each key point will be attributed to specific sections of the referenced documents.

3. If no relevant information is available in the vector database, the response will clearly state that the information is not present. General knowledge or external sources will not be used.

8. Include sufficient details to ensure complete understanding while maintaining clarity.
4. Answers will be precise, structured, and directly relevant to the query, while providing sufficient depth to fully explain concepts. Responses will balance conciseness with comprehensive explanations that ensure complete understanding of the topic. Technical details will be included where necessary, and complex ideas will be broken down into digestible components. Citations will be provided in a reference format, ensuring clarity and credibility.
5. Every claim or piece of information in the response will be accompanied by an inline reference number [1], [2], etc., which corresponds to the citation list at the end. The 'Sources' section will list the exact document name, section, and page number in the following format: [#] Source, Section (page number)

6. If no information is found in the vector database, the response will explicitly state:"No relevant information is available in the vector database."

7. Every response will end with a structured "Sources:" section that lists the exact document names and sections from which the information was retrieved.
8. If you cannot find a direct answer, provide the most relevant information available and indicate limitations.
9. Emphasize visual hierarchy through headings rather than numbered lists.

    """

    # Add error handling for the API call
    try:
        # Retrieve context first
        context_docs = retriever.invoke(enhanced_question)
        context_text = "\n".join([doc.page_content for doc in context_docs])

        # Determine if it's a comparative or exploratory question
        is_comparative = is_comparative_or_exploratory_question(enhanced_question)

        # Prepare the full prompt
        full_prompt = template.format(context=context_text, question=enhanced_question)

        # If it's a comparative or exploratory question, add more context
        if is_comparative and context_text:
            full_prompt += "\n\nNote: This is a comparative or exploratory question. Provide a comprehensive analysis using available context and broader knowledge."

        # Generate response using Gemini
        response = model.generate_content(full_prompt)
        
        return {"answer": response.text}
    except Exception as e:
        return {"answer": f"An error occurred while processing your request. Please try again with a more specific question.", "error": str(e)}, 500


if __name__ == '__main__':
    port = int(os.environ.get('PORT', 3000))

    app.run(debug=False, host='0.0.0.0', port=port)
